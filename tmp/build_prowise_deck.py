from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


OUT = Path(r"C:\Users\T560\Desktop\PRO-WISE\LaTeX_Projects\PRO-WISE_Soutenance_20_slides_clean.pptx")


BLUE = RGBColor(22, 77, 122)
CYAN = RGBColor(22, 156, 166)
INK = RGBColor(28, 38, 50)
MUTED = RGBColor(93, 108, 124)
LIGHT = RGBColor(240, 247, 250)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(66, 146, 110)
ORANGE = RGBColor(232, 142, 67)


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def textbox(slide, x, y, w, h, text="", size=24, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Aptos"
    return box


def add_header(slide, idx, section, title):
    textbox(slide, 0.55, 0.25, 0.8, 0.35, f"{idx:02}", 15, BLUE, True)
    textbox(slide, 1.25, 0.25, 3.1, 0.35, section.upper(), 9, MUTED, False)
    textbox(slide, 0.55, 0.68, 11.6, 0.48, title, 22, INK, True)
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.2), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT
    line.line.fill.background()


def add_footer(slide):
    textbox(slide, 10.8, 7.0, 1.95, 0.25, "PRO-WISE Platform", 8, MUTED, False, PP_ALIGN.RIGHT)


def bullet_list(slide, x, y, w, h, items, size=18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.font.name = "Aptos"
        p.space_after = Pt(7)
    return box


def card(slide, x, y, w, h, title, body, accent=CYAN):
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(218, 229, 235)
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    textbox(slide, x + 0.22, y + 0.16, w - 0.35, 0.35, title, 14, INK, True)
    textbox(slide, x + 0.22, y + 0.65, w - 0.35, h - 0.8, body, 11.5, MUTED)


def placeholder(slide, x, y, w, h, label):
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(246, 250, 252)
    shape.line.color.rgb = RGBColor(180, 205, 218)
    shape.line.dash_style = 4
    textbox(slide, x + 0.15, y + h / 2 - 0.22, w - 0.3, 0.45, label, 16, BLUE, True, PP_ALIGN.CENTER)


def title_slide(prs):
    s = blank_slide(prs)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(246, 250, 252)
    textbox(s, 0.75, 0.65, 11.8, 0.6, "PRO-WISE", 40, BLUE, True, PP_ALIGN.CENTER)
    textbox(s, 1.3, 1.35, 10.7, 0.6, "QR-Based Product Assistance System", 24, INK, True, PP_ALIGN.CENTER)
    textbox(s, 1.45, 2.25, 10.4, 0.7, "Graduation Project Defense - Bachelor's Degree in Computer Science", 17, MUTED, False, PP_ALIGN.CENTER)
    card(s, 1.4, 3.25, 4.95, 1.0, "Presented by", "Ahmed Abderrahmen Jatlaoui", BLUE)
    card(s, 6.95, 3.25, 4.95, 1.0, "Academic year", "2025 / 2026", CYAN)
    textbox(s, 1.3, 4.95, 10.7, 0.6, "Supervised by Mr. Abdelweheb Gueddess and Mr. Mohamed Ben Ameur", 14, MUTED, False, PP_ALIGN.CENTER)
    textbox(s, 1.3, 5.55, 10.7, 0.45, "Host Company: Accelerant", 15, INK, True, PP_ALIGN.CENTER)


slides = [
    ("Introduction & Industrial Context", "Presentation Plan", ["Introduction & Industrial Context", "Problem & Proposed Solution", "Project Objectives & Requirements", "Agile Sprints & Conception Design", "Technical Realization & Prototype Demo", "Conclusion & Perspectives"]),
    ("Introduction & Industrial Context", "General Introduction", ["Paper manuals are lost, outdated, or difficult to search.", "Support teams repeatedly handle simple diagnostic questions.", "QR routing connects the physical product to the right digital support space instantly."]),
    ("Introduction & Industrial Context", "Industrial Context", ["Manufacturers need scalable after-sales support.", "Customers expect immediate self-service assistance.", "Brands need centralized control over product manuals, media, feedback, and maintenance requests."]),
    ("Problem & Proposed Solution", "Problem Statement", ["Support content is scattered across PDFs, videos, and portals.", "Traditional portals do not know the exact product model.", "Simple issues often become costly support calls or field visits."]),
    ("Problem & Proposed Solution", "Proposed Solution", ["Unique QR code per product.", "Web and mobile dashboards with product-specific guides.", "AI-assisted troubleshooting, semantic search, recommendations, and technician escalation."]),
    ("Project Objectives & Requirements", "Project Objectives", ["Digitize product assistance and reduce dependence on paper manuals.", "Provide secure multi-tenant access for brands and administrators.", "Deliver product search, guide management, feedback, notifications, and maintenance workflows."]),
    ("Project Objectives & Requirements", "Requirements", ["Functional: authentication, products, guides, QR, feedback, AI search, technician requests.", "Non-functional: security, reliability, performance, usability, scalability.", "Actors: visitor, user, technician, company admin, super admin."]),
    ("Agile Sprints & Conception Design", "Agile Scrum Methodology", ["The project was divided into five releases.", "Each sprint focused on a coherent functional module.", "UML diagrams and UI validation were used to stabilize the conception before implementation."]),
    ("Agile Sprints & Conception Design", "Product Backlog & Sprint Planning", ["Sprint 1: authentication and user management.", "Sprint 2: catalog, guides, QR, and content approval.", "Sprint 3: feedback, notifications, analytics, and monitoring.", "Sprint 4: AI, RAG, semantic search, and recommendations.", "Sprint 5: maintenance requests and technician workflow."]),
    ("Agile Sprints & Conception Design", "Global Use Case Diagram", ["Use this slide for the global use case diagram already prepared in the report."]),
    ("Agile Sprints & Conception Design", "Global Architecture", ["React web client and Expo mobile client.", "Sails.js REST API with MongoDB persistence.", "Gemini services for embeddings, semantic search, and AI assistance."]),
    ("Agile Sprints & Conception Design", "UML by Sprint", ["Keep your five sprint diagrams here or split them visually into five labeled blocks."]),
    ("Agile Sprints & Conception Design", "RAG and Intelligent Assistance Design", ["Product documents and guide content are indexed as searchable knowledge.", "The system retrieves relevant product context before generating assistance.", "Semantic matching improves troubleshooting and component compatibility suggestions."]),
    ("Technical Realization & Prototype Demo", "Technical Stack", ["Frontend: React, Vite, TypeScript, Redux Toolkit.", "Mobile: React Native, Expo, QR scanner.", "Backend: Node.js, Sails.js, MongoDB, JWT, Gemini API.", "Tools: VS Code, Git, Draw.io, Vercel, Render."]),
    ("Technical Realization & Prototype Demo", "Implemented Modules", ["Authentication and role-based access control.", "Product catalog, guide editor, QR generation, and approval workflow.", "Feedback, notifications, analytics, AI search, and technician maintenance requests."]),
    ("Technical Realization & Prototype Demo", "Prototype Demo Flow", ["Scan QR code.", "Open product detail dashboard.", "Consult guides, manuals, videos, and feedback.", "Use semantic search or AI assistance.", "Escalate unresolved issues to a technician."]),
    ("Conclusion & Perspectives", "Conclusion", ["PRO-WISE centralizes product assistance in a single digital platform.", "The project delivers web, mobile, backend, QR, AI, and maintenance workflows.", "The platform reduces support friction for customers and operational load for manufacturers."]),
    ("Conclusion & Perspectives", "Perspectives", ["Improve mobile offline access for manuals and guides.", "Add richer AI moderation and multilingual support.", "Enhance analytics dashboards and predictive maintenance.", "Extend integrations with manufacturer ERP/CRM systems."]),
    ("Conclusion & Perspectives", "Thank You", ["Thank you for your attention.", "Questions?"]),
]


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs)

    for idx, (section, title, items) in enumerate(slides, start=2):
        slide = blank_slide(prs)
        add_header(slide, idx, section, title)
        add_footer(slide)

        if idx == 2:
            for n, item in enumerate(items):
                card(slide, 0.8 + (n % 2) * 6.1, 1.75 + (n // 2) * 1.45, 5.6, 0.95, f"{n + 1:02}", item, [BLUE, CYAN, GREEN, ORANGE, BLUE, CYAN][n])
        elif idx in (11, 13):
            placeholder(slide, 0.8, 1.65, 11.75, 4.95, "Insert prepared conception diagram / sprint visuals")
            bullet_list(slide, 0.95, 6.0, 11.2, 0.65, items, 12)
        elif idx == 14:
            placeholder(slide, 0.75, 1.55, 5.65, 4.7, "Insert RAG architecture image")
            bullet_list(slide, 6.75, 1.75, 5.65, 4.2, items, 16)
        elif idx == 17:
            placeholder(slide, 0.8, 1.55, 5.7, 4.8, "Insert prototype screenshots")
            bullet_list(slide, 6.9, 1.75, 5.1, 4.4, items, 17)
        elif idx == 20:
            textbox(slide, 1.2, 2.4, 10.9, 0.7, "Thank you for your attention", 34, BLUE, True, PP_ALIGN.CENTER)
            textbox(slide, 1.2, 3.25, 10.9, 0.6, "Questions?", 26, INK, True, PP_ALIGN.CENTER)
        else:
            bullet_list(slide, 1.0, 1.75, 11.2, 4.9, items, 20 if len(items) <= 3 else 17)

    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    main()
